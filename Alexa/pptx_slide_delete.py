import argparse, os, sys
from pathlib import Path
from pptx import Presentation

def _delete_slide(presentation,  index):
    # delete any slide from the index
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
    return True

def rename_file(folder_path,filename):
    # remove the space from the filename
    old_file_path = os.path.join(folder_path,filename)
    filename = filename.replace(' ','_')
    new_file_path = os.path.join(folder_path,filename)
    os.rename(old_file_path, new_file_path)
    return new_file_path

def read_slide(file_path):
    try:
        # open the pptx file and create the pptx object
        sys.stdout.write('the file '+file_path+' is processing....\n')
        prs = Presentation(file_path)
        # remove wrong keyword
        wrong_word = 'wrong'
        for slide in prs.slides:
            index = prs.slides.index(slide)
            status = True
            # Get shape
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if wrong_word in run.text.lower():
                            status = False
                            break;
                    if not status:
                        break;
                if not status:
                        break;
            if not status:
                # if wrong word found then delete the slide
                _delete_slide(prs,index)
        prs.save(file_path)
        sys.stdout.write('Done...\n\n')
    except Exception as e:
        raise Exception(e)

if __name__== "__main__":

    parser = argparse.ArgumentParser()

    parser.add_argument('--fp', action='store', dest='file_path',
                        help='Give a proper file path')
    parser.add_argument('--fd', action='store', dest='folder_path',
                        help='Give a proper folder path')

    results = parser.parse_args()
    file_path = results.file_path
    folder_path = results.folder_path

    if file_path:
        filename = Path(file_path).name
        if '.pptx' in filename:
            read_slide(rename_file(folder_path,filename))
    elif folder_path:
        dirpath = Path(folder_path)
        try:
            for x in dirpath.iterdir():
                filename = Path(x).name
                new_file_path = ''
                if x.is_file() and '.pptx' in filename:
                    read_slide(rename_file(folder_path,filename))
        except Exception as e:
            print(e)
            sys.stderr.write("Please provide a folder path\n")
    else:
        sys.stderr.write("Please provide one option or use help option -h\n")





    